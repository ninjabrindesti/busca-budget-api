import os
import re
import uuid
from copy import deepcopy

import requests
from pptx import Presentation

PLACEHOLDER_PATTERN = re.compile(r"{{\s*([\w\.]+)\s*}}")


def _replace_in_text_frame(text_frame, data: dict):
    for paragraph in text_frame.paragraphs:
        original_text = "".join(run.text for run in paragraph.runs)

        if not original_text:
            continue

        def repl(match):
            key = match.group(1)
            value = data.get(key)

            if value is None or value == "":
                return match.group(0)

            return str(value)

        new_text = PLACEHOLDER_PATTERN.sub(repl, original_text)

        if new_text != original_text and len(paragraph.runs) > 0:
            paragraph.runs[0].text = new_text

            for run in paragraph.runs[1:]:
                run.text = ""


def replace_text_placeholders(prs: Presentation, data: dict):
    for slide in prs.slides:
        replace_text_placeholders_on_slide(slide, data)


def replace_text_placeholders_on_slide(slide, data: dict):
    for shape in slide.shapes:
        if getattr(shape, "has_table", False):
            for row in shape.table.rows:
                for cell in row.cells:
                    if cell.text_frame is not None:
                        _replace_in_text_frame(cell.text_frame, data)
            continue

        if hasattr(shape, "text_frame") and shape.text_frame is not None:
            _replace_in_text_frame(shape.text_frame, data)


def _download_image(url: str) -> str | None:
    if not url:
        return None

    response = requests.get(url, timeout=20)
    response.raise_for_status()

    content_type = response.headers.get("content-type", "").lower()

    extension = ".png"
    if "jpeg" in content_type or "jpg" in content_type:
        extension = ".jpg"

    file_path = f"/tmp/{uuid.uuid4().hex}{extension}"

    with open(file_path, "wb") as f:
        f.write(response.content)

    return file_path


def replace_named_images(prs: Presentation, data: dict):
    for slide in prs.slides:
        replace_named_images_on_slide(slide, data)


def replace_named_images_on_slide(slide, data: dict):
    image_mappings = {
        "seller_image": data.get("seller_image_url"),
        "item_image": data.get("item_image_url"),
    }

    shapes_to_replace = []

    for shape in slide.shapes:
        shape_name = getattr(shape, "name", "")

        if shape_name in image_mappings and image_mappings[shape_name]:
            shapes_to_replace.append((shape, image_mappings[shape_name]))

    for shape, image_url in shapes_to_replace:
        image_path = _download_image(image_url)
        if not image_path:
            continue

        left = shape.left
        top = shape.top
        width = shape.width
        height = shape.height

        sp = shape._element
        sp.getparent().remove(sp)

        slide.shapes.add_picture(image_path, left, top, width=width, height=height)

        if os.path.exists(image_path):
            os.remove(image_path)


def duplicate_slide(prs, slide_index: int):
    source_slide = prs.slides[slide_index]
    blank_layout = prs.slide_layouts[6]
    new_slide = prs.slides.add_slide(blank_layout)

    for shape in source_slide.shapes:
        new_el = deepcopy(shape._element)
        new_slide.shapes._spTree.insert_element_before(new_el, "p:extLst")

    return new_slide
