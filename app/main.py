"""
main.py — Gerador de Proposta Ninja Brindes
Sprint 1: contrato único de placeholders + payload legado/novo
Sprint 2: motor central de replace (text frames, tabelas, multi-run)
Sprint 3: schema de pagamentos flexível, paginação do slide de resumo (>5 itens),
          preenchimento do slide de pagamentos e remoção de linhas vazias.
Sprint 4: merge automático dos slides de orçamento + pagamento quando itens <= ITEMS_MERGE_THRESHOLD.
          Quando há poucos itens, os slides 10 e 11 são fundidos em um único slide (slide 10),
          com a tabela de itens no topo e o bloco de pagamentos embutido abaixo — sem alterar
          o backend nem o payload da API.
"""

import io
import math
import os
import re as _re
import traceback
import uuid
import zipfile as _zipfile
from contextlib import asynccontextmanager
from copy import deepcopy as _deepcopy
from typing import List, Optional, Iterable
from urllib.parse import urlparse

import httpx
from lxml import etree as _etree

from fastapi import FastAPI, HTTPException
from fastapi.responses import StreamingResponse
from pydantic import BaseModel, ConfigDict, Field
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from app.services.pptx_generator import replace_named_images_on_slide
from app.services.slide_duplicator import duplicate_slide_in_pptx


# ---------------------------------------------------------------------------
# Constantes
# ---------------------------------------------------------------------------

TEMPLATE_PATH = "templates/template_ninja.pptx"
TEMPLATE_URL = os.getenv("TEMPLATE_URL")

COVER_SLIDE_INDEX   = 0
ITEM_SLIDE_INDEX    = 8   # Slide 9  (0-based) — template de item
SUMMARY_SLIDE_INDEX = 9   # Slide 10 (0-based) — tabela de itens
PAYMENT_SLIDE_INDEX = 10  # Slide 11 (0-based) — resumo de pagamentos

ITEMS_PER_SUMMARY_SLIDE = 5   # Máximo de itens por slide de resumo
ITEMS_MERGE_THRESHOLD   = 3   # Itens <= este valor → slides 10 e 11 fundidos em um único slide


# ---------------------------------------------------------------------------
# Template bootstrap
# ---------------------------------------------------------------------------

def _is_valid_http_url(url: Optional[str]) -> bool:
    if not url or not isinstance(url, str):
        return False
    url = url.strip()
    if " " in url:
        return False
    try:
        parsed = urlparse(url)
    except Exception:
        return False
    return parsed.scheme in ("http", "https") and bool(parsed.netloc)


@asynccontextmanager
async def lifespan(app: FastAPI):
    os.makedirs("templates", exist_ok=True)
    if not os.path.exists(TEMPLATE_PATH):
        if not _is_valid_http_url(TEMPLATE_URL):
            raise RuntimeError("Template ausente e TEMPLATE_URL inválida.")
        try:
            with httpx.Client(timeout=30, follow_redirects=True) as client:
                response = client.get(TEMPLATE_URL.strip())
                response.raise_for_status()
        except Exception as exc:
            raise RuntimeError(f"Falha ao baixar template: {exc}") from exc
        with open(TEMPLATE_PATH, "wb") as f:
            f.write(response.content)
    yield


app = FastAPI(title="Ninja Brindes - Gerador de Proposta", lifespan=lifespan)


# ---------------------------------------------------------------------------
# Schemas
# ---------------------------------------------------------------------------

class PaymentEntry(BaseModel):
    """
    Representa uma linha de pagamento (entrada ou parcela).

    Convenção de uso na lista `proposal.payments`:
      COM entrada:
        [0] → Entrada:        label="Entrada", method="PIX", date="27/04/2026", value=4000.0
        [1] → Header parcela: label="Parcelamento", method="Boleto", plan="Boleto 4x"
        [2..N] → Parcelas:    date="08/05/2026", value=1000.0

      SEM entrada:
        [0] → Header parcela: label="Parcelamento", method="Boleto", plan="Boleto 4x"
        [1..N] → Parcelas:    date="08/05/2026", value=1000.0

    A detecção é feita por conteúdo (não por posição) via _is_entry_row().
    Campos vazios/None → célula em branco. Linhas sem nenhum dado são removidas do slide.
    """
    model_config = ConfigDict(extra="allow")
    label:  Optional[str]   = ""    # {{entry}} / {{installments}}
    method: Optional[str]   = ""    # {{entry_method}} / {{installments_method}}
    plan:   Optional[str]   = ""    # {{installments_plan}}
    date:   Optional[str]   = ""    # {{payments_schedule}}
    value:  Optional[float] = None  # {{entry_value}} / {{installments_value}}


class Proposal(BaseModel):
    model_config = ConfigDict(extra="allow")
    proposal_number:     str
    client_name:         str
    company_name:        Optional[str] = ""
    seller_name:         Optional[str] = ""
    seller_phone:        Optional[str] = ""
    seller_email:        Optional[str] = ""
    seller_description:  Optional[str] = ""
    seller_image_url:    Optional[str] = None
    payment_method:      Optional[str] = ""
    payment_term:        Optional[str] = ""
    delivery_date:       Optional[str] = ""
    notes:               Optional[str] = ""
    obs_cnpj:            Optional[str] = ""
    cover_cnpj:          Optional[str] = ""
    cover_corporate_name: Optional[str] = ""

    # Dados de pagamento flexíveis para o slide 11
    payments: Optional[List[PaymentEntry]] = Field(default_factory=list)

    # Campos legados / retrocompat
    summary_payment_method: Optional[str] = ""
    summary_delivery_date:  Optional[str] = ""


class Item(BaseModel):
    model_config = ConfigDict(extra="allow")
    item_index:       int
    item_name:        str             = Field(max_length=120)
    item_subtitle:    Optional[str]   = Field(default="", max_length=120)
    item_description: Optional[str]   = Field(default="", max_length=500)
    item_code:        Optional[str]   = Field(default="", max_length=60)
    quantity:         int             = 1
    unit_price:       float           = 0.0
    item_image_url:   Optional[str]   = ""


class Section(BaseModel):
    model_config = ConfigDict(extra="allow")
    section_index:  int
    freight_value:  Optional[float] = None
    freight_label:  Optional[str]   = ""
    items:          List[Item]      = Field(min_length=1, max_length=50)


class GenerateRequest(BaseModel):
    model_config = ConfigDict(extra="allow")
    proposal: Proposal
    sections: List[Section] = Field(min_length=1)


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _format_currency(value: float) -> str:
    return f"R$ {value:,.2f}".replace(",", "X").replace(".", ",").replace("X", ".")


def _calculate_section_total(section: Section) -> float:
    return sum(item.quantity * item.unit_price for item in section.items)


def _calculate_grand_total(section: Section) -> float:
    """grand_total = subtotal dos itens + frete"""
    return _calculate_section_total(section) + (section.freight_value or 0.0)


def _s(value) -> str:
    return "" if value is None else str(value)


# ---------------------------------------------------------------------------
# Contrato de placeholders
# ---------------------------------------------------------------------------

def _build_global_data(proposal: Proposal) -> dict:
    company_name   = _s(getattr(proposal, "company_name", ""))
    payment_method = _s(getattr(proposal, "payment_method", ""))
    payment_term   = _s(getattr(proposal, "payment_term", ""))
    delivery_date  = _s(getattr(proposal, "delivery_date", ""))
    notes          = _s(getattr(proposal, "notes", ""))
    obs_cnpj       = _s(getattr(proposal, "obs_cnpj", ""))

    in_summary_pm      = _s(getattr(proposal, "summary_payment_method", ""))
    in_summary_dd      = _s(getattr(proposal, "summary_delivery_date", ""))
    in_cover_pn        = _s(getattr(proposal, "cover_proposal_number", ""))
    in_cover_co        = _s(getattr(proposal, "cover_company", ""))
    in_cover_cl        = _s(getattr(proposal, "cover_client", ""))
    in_cover_obs       = _s(getattr(proposal, "cover_obs_cnpj", ""))
    in_cover_cnpj      = _s(getattr(proposal, "cover_cnpj", ""))
    in_cover_corp_name = _s(getattr(proposal, "cover_corporate_name", ""))

    return {
        # legados (retrocompat)
        "proposal_number":    _s(proposal.proposal_number),
        "client_name":        _s(proposal.client_name),
        "company_name":       company_name,
        "seller_name":        _s(getattr(proposal, "seller_name", "")),
        "seller_phone":       _s(getattr(proposal, "seller_phone", "")),
        "seller_email":       _s(getattr(proposal, "seller_email", "")),
        "seller_description": _s(getattr(proposal, "seller_description", "")),
        "seller_image_url":   _s(getattr(proposal, "seller_image_url", "")),
        "payment_method":     payment_method,
        "payment_term":       payment_term,
        "delivery_date":      delivery_date,
        "notes":              notes,
        # oficiais
        "summary_payment_method": in_summary_pm or payment_method or payment_term or "",
        "summary_delivery_date":  in_summary_dd or delivery_date or "",
        "cover_proposal_number":  in_cover_pn   or _s(proposal.proposal_number),
        "cover_company":          in_cover_co   or company_name or "",
        "cover_client":           in_cover_cl   or _s(proposal.client_name),
        "cover_obs_cnpj":         in_cover_obs  or obs_cnpj or notes or "",
        "cover_cnpj":             in_cover_cnpj,
        "cover_corporate_name":   in_cover_corp_name,
    }


def _build_data(proposal: Proposal, item: Item, section: Section, display_index: int) -> dict:
    item_total    = item.quantity * item.unit_price
    section_total = _calculate_section_total(section)
    freight_value = section.freight_value or 0.0
    grand_total   = section_total + freight_value
    return {
        **_build_global_data(proposal),
        "item_name":          _s(item.item_name),
        "item_subtitle":      _s(item.item_subtitle),
        "item_index":         str(item.item_index),
        "item_display_index": str(display_index),
        "item_description":   _s(item.item_description),
        "item_code":          _s(item.item_code),
        "quantity":           str(item.quantity),
        "unit_price":         _format_currency(item.unit_price),
        "item_total":         _format_currency(item_total),
        "item_image_url":     _s(item.item_image_url),
        "section_total":      _format_currency(section_total),
        "freight":            _format_currency(freight_value),
        "section_freight":    _s(section.freight_label),
        "grand_total":        _format_currency(grand_total),
        "total":              _format_currency(grand_total),
        "freight_label":      _s(section.freight_label),
    }


def _build_summary_data(proposal: Proposal, section: Section) -> dict:
    section_total = _calculate_section_total(section)
    freight_value = section.freight_value or 0.0
    grand_total   = section_total + freight_value
    return {
        **_build_global_data(proposal),
        "section_total":   _format_currency(section_total),
        "freight":         _format_currency(freight_value),
        "section_freight": _s(section.freight_label),
        "grand_total":     _format_currency(grand_total),
        "total":           _format_currency(grand_total),
        "freight_label":   _s(section.freight_label),
    }


# ---------------------------------------------------------------------------
# Detecção de entrada por conteúdo (fix: propostas sem entrada)
# ---------------------------------------------------------------------------

def _is_entry_row(p: PaymentEntry) -> bool:
    """
    Detecta por CONTEÚDO se um PaymentEntry representa uma entrada (down payment).
    """
    if p is None:
        return False
    label = (p.label or "").strip().lower()
    if "entrada" in label:
        return True
    if p.plan and p.plan.strip():
        return False
    if p.value is not None and p.date and p.date.strip():
        return True
    return False


def _parse_payments(payments: list) -> tuple:
    """
    Separa a lista de payments em (entry, installment_header, installment_rows).
    """
    if not payments:
        return None, None, []

    if _is_entry_row(payments[0]):
        entry              = payments[0]
        installment_header = payments[1] if len(payments) > 1 else None
        installment_rows   = payments[2:] if len(payments) > 2 else []
    else:
        entry              = None
        installment_header = payments[0]
        installment_rows   = payments[1:] if len(payments) > 1 else []

    return entry, installment_header, installment_rows


def _build_payment_data(proposal: Proposal, section: Section) -> tuple[dict, list]:
    """
    Monta o dicionário de placeholders para o slide de pagamentos.
    """
    payments = proposal.payments or []
    entry, installment_header, installment_rows = _parse_payments(payments)

    data = {**_build_summary_data(proposal, section)}

    data["entry"]        = _s(entry.label  if entry else "")
    data["entry_method"] = _s(entry.method if entry else "")
    data["entry_value"]  = (
        _format_currency(entry.value) if entry and entry.value is not None else ""
    )
    data["payments_schedule"] = _s(entry.date if entry else "")

    data["installments"]        = _s(installment_header.label  if installment_header else "")
    data["installments_method"] = _s(installment_header.method if installment_header else "")
    data["installments_plan"]   = _s(installment_header.plan   if installment_header else "")

    data["installments_value"] = (
        _format_currency(installment_rows[0].value)
        if installment_rows and installment_rows[0].value is not None
        else ""
    )

    return data, installment_rows


# ===========================================================================
# Motor central de replace
# ===========================================================================

_PLACEHOLDER_CACHE: dict[str, _re.Pattern] = {}


def _pattern_for(key: str) -> _re.Pattern:
    pat = _PLACEHOLDER_CACHE.get(key)
    if pat is None:
        pat = _re.compile(r"\{\{\s*" + _re.escape(key) + r"\s*\}\}")
        _PLACEHOLDER_CACHE[key] = pat
    return pat


def _apply_to_paragraph(paragraph, replacements: dict) -> bool:
    runs = paragraph.runs
    if not runs:
        return False
    original = "".join(r.text or "" for r in runs)
    if "{{" not in original:
        return False
    new_text = original
    for key, value in replacements.items():
        pat = _pattern_for(key)
        if pat.search(new_text):
            new_text = pat.sub(_s(value), new_text)
    if new_text == original:
        return False
    runs[0].text = new_text
    for r in runs[1:]:
        r.text = ""
    return True


def _is_group_shape(shape) -> bool:
    try:
        return shape.shape_type == MSO_SHAPE_TYPE.GROUP
    except (AttributeError, ValueError):
        return False


def _iter_text_frames(shapes) -> Iterable:
    for shape in shapes:
        if _is_group_shape(shape):
            yield from _iter_text_frames(shape.shapes)
            continue
        if getattr(shape, "has_table", False):
            for row in shape.table.rows:
                for cell in row.cells:
                    if cell.text_frame is not None:
                        yield cell.text_frame
            continue
        if getattr(shape, "has_text_frame", False) and shape.text_frame is not None:
            yield shape.text_frame


def replace_placeholders_everywhere(slide, replacements: dict) -> int:
    changed = 0
    for tf in _iter_text_frames(slide.shapes):
        for paragraph in tf.paragraphs:
            if _apply_to_paragraph(paragraph, replacements):
                changed += 1
    return changed


# ---------------------------------------------------------------------------
# Detecção de modo merge
# ---------------------------------------------------------------------------

def _should_merge_slides(sections: List[Section]) -> bool:
    """
    Retorna True se a proposta tem itens suficientemente poucos para
    fundir os slides de orçamento e pagamento em um único slide.

    Critério: total de itens em todas as seções <= ITEMS_MERGE_THRESHOLD.
    Propostas com múltiplas seções NÃO fazem merge (cada seção tem seu próprio
    par resumo+pagamento).
    """
    if len(sections) != 1:
        return False
    return len(sections[0].items) <= ITEMS_MERGE_THRESHOLD


# ---------------------------------------------------------------------------
# Endpoint principal
# ---------------------------------------------------------------------------

@app.get("/")
def health():
    return {"status": "ok"}


@app.post("/generate")
def generate_proposal(payload: GenerateRequest):
    try:
        if not os.path.exists(TEMPLATE_PATH):
            raise HTTPException(status_code=500, detail="Template não encontrado.")

        with open(TEMPLATE_PATH, "rb") as f:
            pptx_bytes = f.read()

        merge_mode = _should_merge_slides(payload.sections)

        # -------------------------------------------------------------------
        # 1. Calcular quantidade de slides necessários
        # -------------------------------------------------------------------
        total_item_slides = sum(len(s.items) for s in payload.sections)
        if total_item_slides <= 0:
            raise HTTPException(status_code=400, detail="Nenhum item enviado.")

        # Slides de resumo por seção: 1 slide a cada 5 itens
        summary_slide_counts = [
            math.ceil(len(s.items) / ITEMS_PER_SUMMARY_SLIDE)
            for s in payload.sections
        ]
        total_summary_slides = sum(summary_slide_counts)

        # -------------------------------------------------------------------
        # 2. Duplicar slides de ITEM (template índice 8)
        # -------------------------------------------------------------------
        for _ in range(max(total_item_slides - 1, 0)):
            pptx_bytes = duplicate_slide_in_pptx(
                input_bytes=pptx_bytes,
                source_slide_index=ITEM_SLIDE_INDEX,
                copies=1,
                insert_after_index=ITEM_SLIDE_INDEX,
            )

        # -------------------------------------------------------------------
        # 3. Duplicar slides de RESUMO (template índice 9, agora deslocado)
        # -------------------------------------------------------------------
        summary_template_index = ITEM_SLIDE_INDEX + total_item_slides
        for _ in range(max(total_summary_slides - 1, 0)):
            pptx_bytes = duplicate_slide_in_pptx(
                input_bytes=pptx_bytes,
                source_slide_index=summary_template_index,
                copies=1,
                insert_after_index=summary_template_index,
            )

        # -------------------------------------------------------------------
        # 4. Reordenar slides
        # -------------------------------------------------------------------
        pptx_bytes = _reorder_slides(
            pptx_bytes=pptx_bytes,
            sections=payload.sections,
            total_item_slides=total_item_slides,
            summary_slide_counts=summary_slide_counts,
        )

        # -------------------------------------------------------------------
        # 5. Substituir placeholders
        # -------------------------------------------------------------------
        prs = Presentation(io.BytesIO(pptx_bytes))
        global_data = _build_global_data(payload.proposal)

        # Capa
        replace_placeholders_everywhere(prs.slides[COVER_SLIDE_INDEX], global_data)
        replace_named_images_on_slide(prs.slides[COVER_SLIDE_INDEX], global_data)

        # Vendedor (slide 4, índice 3 — fixo)
        replace_placeholders_everywhere(prs.slides[3], global_data)
        replace_named_images_on_slide(prs.slides[3], global_data)

        # Itens e resumos por seção
        slide_cursor = ITEM_SLIDE_INDEX
        for sec_i, section in enumerate(payload.sections):

            # Slides de item
            for display_index, item in enumerate(section.items, start=1):
                if slide_cursor >= len(prs.slides):
                    raise HTTPException(
                        status_code=500,
                        detail=f"Slide de item ausente. Índice: {slide_cursor}",
                    )
                slide = prs.slides[slide_cursor]
                data = _build_data(payload.proposal, item, section, display_index)
                replace_placeholders_everywhere(slide, data)
                replace_named_images_on_slide(slide, data)
                slide_cursor += 1

            # Slides de resumo (1 por cada 5 itens)
            n_summary = summary_slide_counts[sec_i]
            for page in range(n_summary):
                if slide_cursor >= len(prs.slides):
                    raise HTTPException(
                        status_code=500,
                        detail=f"Slide de resumo ausente. Índice: {slide_cursor}",
                    )
                summary_slide = prs.slides[slide_cursor]
                summary_data  = _build_summary_data(payload.proposal, section)

                # Fatia de itens para esta página
                start = page * ITEMS_PER_SUMMARY_SLIDE
                end   = start + ITEMS_PER_SUMMARY_SLIDE
                section_chunk = section.model_copy(
                    update={"items": section.items[start:end]}
                )

                _expand_summary_table_rows(summary_slide, section_chunk)
                replace_placeholders_everywhere(summary_slide, summary_data)
                _replace_summary_table_rows(summary_slide, section_chunk, index_offset=start)
                replace_named_images_on_slide(summary_slide, summary_data)

                # ── MERGE MODE ──────────────────────────────────────────────
                # Quando há poucos itens (seção única, ≤ ITEMS_MERGE_THRESHOLD),
                # incorporamos o conteúdo do slide de pagamentos DIRETAMENTE
                # neste slide de resumo e pulamos o slide de pagamentos original.
                if merge_mode and page == 0:
                    payment_data, installment_rows = _build_payment_data(
                        payload.proposal, section
                    )
                    _merge_payment_into_summary_slide(
                        summary_slide, payload.proposal, installment_rows, payment_data
                    )
                    replace_placeholders_everywhere(summary_slide, payment_data)

                slide_cursor += 1

        # -------------------------------------------------------------------
        # Slide de pagamentos (slide 11)
        # Em merge_mode ele é exibido mas sem conteúdo de pagamento — removemos
        # as formas de pagamento para evitar placeholders soltos, mas mantemos
        # o slide na apresentação para não quebrar o índice dos slides finais.
        # Melhor ainda: simplesmente o pulamos e removemos do output.
        # -------------------------------------------------------------------
        if slide_cursor < len(prs.slides):
            payment_slide = prs.slides[slide_cursor]

            if merge_mode:
                # Remove o slide de pagamentos do output (já foi embutido no resumo)
                _delete_slide(prs, slide_cursor)
                # slide_cursor não avança — o próximo slide (contato) já está nesta posição
            else:
                payment_data, installment_rows = _build_payment_data(
                    payload.proposal,
                    payload.sections[-1],
                )
                _remove_empty_payment_rows(payment_slide, payload.proposal)
                _expand_installment_rows(payment_slide, installment_rows, payment_data)
                replace_placeholders_everywhere(payment_slide, payment_data)
                replace_named_images_on_slide(payment_slide, payment_data)
                slide_cursor += 1

        # Slide final (contato)
        last_slide = prs.slides[-1]
        replace_placeholders_everywhere(last_slide, global_data)
        replace_named_images_on_slide(last_slide, global_data)

        out_buf = io.BytesIO()
        prs.save(out_buf)
        out_buf.seek(0)

        filename = f"proposta_{uuid.uuid4().hex[:8]}.pptx"
        return StreamingResponse(
            out_buf,
            media_type=(
                "application/vnd.openxmlformats-officedocument"
                ".presentationml.presentation"
            ),
            headers={"Content-Disposition": f'attachment; filename="{filename}"'},
        )

    except HTTPException:
        raise
    except Exception as e:
        traceback.print_exc()
        raise HTTPException(
            status_code=500,
            detail={"error": "render_failed", "type": type(e).__name__, "message": str(e)},
        )


# ---------------------------------------------------------------------------
# Reordenação dos slides
# ---------------------------------------------------------------------------

_NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"


def _reorder_slides(pptx_bytes, sections, total_item_slides, summary_slide_counts):
    total_summary_slides = sum(summary_slide_counts)
    buf: dict[str, bytes] = {}

    with _zipfile.ZipFile(io.BytesIO(pptx_bytes), "r") as zf:
        prs_root   = _etree.fromstring(zf.read("ppt/presentation.xml"))
        sld_id_lst = prs_root.find(f".//{{{_NS_P}}}sldIdLst")
        sld_id_els = list(sld_id_lst)

        fixed_before   = sld_id_els[:ITEM_SLIDE_INDEX]
        item_slides    = sld_id_els[ITEM_SLIDE_INDEX: ITEM_SLIDE_INDEX + total_item_slides]
        summary_slides = sld_id_els[
            ITEM_SLIDE_INDEX + total_item_slides:
            ITEM_SLIDE_INDEX + total_item_slides + total_summary_slides
        ]
        fixed_after = sld_id_els[
            ITEM_SLIDE_INDEX + total_item_slides + total_summary_slides:
        ]

        new_order = list(fixed_before)
        item_ptr = summary_ptr = 0
        for sec_i, section in enumerate(sections):
            for _ in range(len(section.items)):
                new_order.append(item_slides[item_ptr])
                item_ptr += 1
            for _ in range(summary_slide_counts[sec_i]):
                new_order.append(summary_slides[summary_ptr])
                summary_ptr += 1
        new_order.extend(fixed_after)

        for el in list(sld_id_lst):
            sld_id_lst.remove(el)
        for el in new_order:
            sld_id_lst.append(el)

        buf["ppt/presentation.xml"] = _etree.tostring(
            prs_root, xml_declaration=True, encoding="utf-8", standalone=True,
        )

        out = io.BytesIO()
        with _zipfile.ZipFile(out, "w", _zipfile.ZIP_DEFLATED) as out_zf:
            written = set()
            for name, data in buf.items():
                out_zf.writestr(name, data)
                written.add(name)
            for name in zf.namelist():
                if name not in written:
                    out_zf.writestr(name, zf.read(name))
        return out.getvalue()


# ---------------------------------------------------------------------------
# Deleção de slide via python-pptx (zip-level)
# ---------------------------------------------------------------------------

def _delete_slide(prs: Presentation, slide_index: int):
    """
    Remove um slide da apresentação pelo índice (0-based).
    Opera diretamente no XML interno do python-pptx.
    """
    xml_slides = prs.slides._sldIdLst
    # Lista de rId dos slides na ordem atual
    slide_ids = list(xml_slides)
    if slide_index >= len(slide_ids):
        return

    # rId do slide a remover
    rId = slide_ids[slide_index].get(
        "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
    )

    # Remove da lista de slides
    xml_slides.remove(slide_ids[slide_index])

    # Remove a relationship da presentation
    del prs.part.part_related_by(rId)


# ---------------------------------------------------------------------------
# Merge: copia shapes de pagamento para dentro do slide de resumo
# ---------------------------------------------------------------------------

def _merge_payment_into_summary_slide(
    summary_slide,
    proposal: Proposal,
    installment_rows: list,
    payment_data: dict,
):
    """
    Incorpora o conteúdo do slide de pagamentos (slide 11 do template) no
    slide de resumo (slide 10), copiando todas as shapes de pagamento.

    Estratégia:
    - Abre o template original para obter o slide 11 limpo.
    - Copia cada shape do slide 11 para o slide 10.
    - Aplica _remove_empty_payment_rows e _expand_installment_rows
      diretamente nas shapes copiadas (operando no slide de resumo).

    Obs: esta função modifica summary_slide in-place.
    """
    if not os.path.exists(TEMPLATE_PATH):
        return

    with open(TEMPLATE_PATH, "rb") as f:
        template_bytes = f.read()

    template_prs = Presentation(io.BytesIO(template_bytes))

    # Slide 11 do template (índice 10 = PAYMENT_SLIDE_INDEX)
    if PAYMENT_SLIDE_INDEX >= len(template_prs.slides):
        return

    payment_template_slide = template_prs.slides[PAYMENT_SLIDE_INDEX]

    # Copia todas as shapes do slide de pagamento para o slide de resumo
    for shape in payment_template_slide.shapes:
        new_el = _deepcopy(shape._element)
        summary_slide.shapes._spTree.insert_element_before(new_el, "p:extLst")

    # Agora opera no slide de resumo (que já tem as shapes de pagamento)
    _remove_empty_payment_rows(summary_slide, proposal)
    _expand_installment_rows(summary_slide, installment_rows, payment_data)


# ---------------------------------------------------------------------------
# Tabela de resumo de itens (slide 10)
# ---------------------------------------------------------------------------

_ITEM_ROW_TOKENS = (
    "{{item_", "{{ item_", "{{quantity", "{{ quantity",
    "{{unit_price", "{{ unit_price", "{{item_total", "{{ item_total",
)


def _row_text(row) -> str:
    return " ".join(cell.text for cell in row.cells)


def _is_item_template_row(row) -> bool:
    text = _row_text(row)
    return any(tok in text for tok in _ITEM_ROW_TOKENS)


def _find_first_item_row_index(table) -> Optional[int]:
    for i, row in enumerate(table.rows):
        if _is_item_template_row(row):
            return i
    return None


def _find_summary_table(slide):
    for shape in slide.shapes:
        if not getattr(shape, "has_table", False):
            continue
        if _find_first_item_row_index(shape.table) is not None:
            return shape.table
    return None


def _expand_summary_table_rows(slide, section):
    table = _find_summary_table(slide)
    if table is None:
        return
    template_idx = _find_first_item_row_index(table)
    if template_idx is None:
        return
    template_row_el = list(table.rows)[template_idx]._tr
    parent = template_row_el.getparent()
    # Remove extras
    for row in list(table.rows)[template_idx + 1:]:
        if _is_item_template_row(row):
            parent.remove(row._tr)
    # Duplica para cada item adicional
    anchor = template_row_el
    for _ in section.items[1:]:
        new_row_el = _deepcopy(template_row_el)
        parent.insert(parent.index(anchor) + 1, new_row_el)
        anchor = new_row_el


def _replace_summary_table_rows(slide, section, index_offset: int = 0):
    table = _find_summary_table(slide)
    if table is None:
        return
    first_idx = _find_first_item_row_index(table)
    if first_idx is None:
        return
    rows      = list(table.rows)
    item_rows = rows[first_idx: first_idx + len(section.items)]
    if len(item_rows) < len(section.items):
        return
    for display_index, (row, item) in enumerate(zip(item_rows, section.items), start=1 + index_offset):
        item_total = item.quantity * item.unit_price
        row_data = {
            "item_display_index": str(display_index),
            "item_index":         str(item.item_index),
            "item_name":          _s(item.item_name),
            "item_subtitle":      _s(item.item_subtitle),
            "item_description":   _s(item.item_description),
            "item_code":          _s(item.item_code),
            "quantity":           str(item.quantity),
            "unit_price":         _format_currency(item.unit_price),
            "item_total":         _format_currency(item_total),
        }
        for cell in row.cells:
            if cell.text_frame is None:
                continue
            for paragraph in cell.text_frame.paragraphs:
                _apply_to_paragraph(paragraph, row_data)


# ---------------------------------------------------------------------------
# Slide de pagamentos (slide 11): remoção de linhas vazias + expansão de parcelas
# ---------------------------------------------------------------------------

def _find_payment_table(slide):
    """Retorna a tabela do slide de pagamentos (identifica por placeholders)."""
    for shape in slide.shapes:
        if not getattr(shape, "has_table", False):
            continue
        table = shape.table
        for row in table.rows:
            text = _row_text(row)
            if any(tok in text for tok in (
                "{{entry}}", "{{installments}}", "{{payments_schedule}}"
            )):
                return table
    return None


def _remove_empty_payment_rows(slide, proposal: Proposal):
    """
    Remove linhas da tabela de pagamentos que não têm dados.
    """
    payments = proposal.payments or []
    table    = _find_payment_table(slide)
    if table is None:
        return

    entry, installment_header, _ = _parse_payments(payments)
    has_schedule = len(payments) > 0

    rows_to_remove = []
    for row in table.rows:
        text = _row_text(row)

        if "{{entry}}" in text:
            if not entry or not (entry.label or entry.method or entry.value is not None):
                rows_to_remove.append(row._tr)

        elif "{{installments}}" in text:
            if not installment_header or not (installment_header.label or installment_header.method):
                rows_to_remove.append(row._tr)

        elif "{{payments_schedule}}" in text:
            if not has_schedule:
                rows_to_remove.append(row._tr)

    for tr in rows_to_remove:
        p = tr.getparent()
        if p is not None:
            p.remove(tr)


def _expand_installment_rows(slide, installment_rows: list, payment_data: dict):
    """
    Expande as linhas de parcela no slide de pagamentos.
    """
    table = _find_payment_table(slide)
    if table is None or not installment_rows:
        return

    schedule_rows = [
        row for row in table.rows
        if "{{payments_schedule}}" in _row_text(row)
    ]
    if not schedule_rows:
        return

    installment_template_row = schedule_rows[-1]
    template_el = installment_template_row._tr
    parent      = template_el.getparent()

    anchor = template_el
    for _ in installment_rows[1:]:
        new_row_el = _deepcopy(template_el)
        parent.insert(parent.index(anchor) + 1, new_row_el)
        anchor = new_row_el

    all_schedule_rows = [
        row for row in table.rows
        if "{{payments_schedule}}" in _row_text(row)
    ]
    parcela_rows = all_schedule_rows[len(all_schedule_rows) - len(installment_rows):]

    for row, inst in zip(parcela_rows, installment_rows):
        installments_method = _s(inst.method) or payment_data.get("installments_method", "")
        row_data = {
            "payments_schedule":   _s(inst.date),
            "installments_method": installments_method,
            "installments_value":  (
                _format_currency(inst.value) if inst.value is not None else ""
            ),
            "entry_value":         "",
            "entry_method":        "",
        }
        for cell in row.cells:
            if cell.text_frame is None:
                continue
            for paragraph in cell.text_frame.paragraphs:
                _apply_to_paragraph(paragraph, row_data)
